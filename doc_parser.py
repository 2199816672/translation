"""
文档解析模块 - 支持 TXT/PDF/Word/Excel/PPT/图片 文本提取
"""
import os


SUPPORTED_EXTENSIONS = {
    '.txt': '文本文件',
    '.pdf': 'PDF 文档',
    '.docx': 'Word 文档',
    '.pptx': 'PowerPoint',
    '.xlsx': 'Excel 表格',
    '.png': '图片',
    '.jpg': '图片',
    '.jpeg': '图片',
    '.bmp': '图片',
    '.gif': '图片',
}

IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.bmp', '.gif'}


def get_file_type(path):
    ext = os.path.splitext(path)[1].lower()
    return SUPPORTED_EXTENSIONS.get(ext, None)


def is_supported(path):
    return get_file_type(path) is not None


def parse_file(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == '.txt':
        return _parse_txt(path)
    elif ext == '.pdf':
        return _parse_pdf(path)
    elif ext == '.docx':
        return _parse_docx(path)
    elif ext == '.pptx':
        return _parse_pptx(path)
    elif ext == '.xlsx':
        return _parse_xlsx(path)
    elif ext in IMAGE_EXTENSIONS:
        return _parse_image(path)
    else:
        raise ValueError(f"不支持的文件格式: {ext}")


def _parse_txt(path):
    for enc in ('utf-8', 'gbk', 'gb2312', 'latin-1'):
        try:
            with open(path, 'r', encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    raise RuntimeError("无法读取文件，请检查文件编码")


def _parse_pdf(path):
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text_parts.append(t)
        if text_parts:
            return '\n\n'.join(text_parts)
    except Exception:
        pass
    try:
        from pdfminer.high_level import extract_text
        return extract_text(path)
    except ImportError:
        raise RuntimeError("请安装 pdfplumber 或 pdfminer.six: pip install pdfplumber")


def _parse_docx(path):
    try:
        from docx import Document
        doc = Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return '\n\n'.join(paragraphs)
    except ImportError:
        raise RuntimeError("请安装 python-docx: pip install python-docx")


def _parse_pptx(path):
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            text_parts.append(t)
        return '\n\n'.join(text_parts)
    except ImportError:
        raise RuntimeError("请安装 python-pptx: pip install python-pptx")


def _parse_xlsx(path):
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        text_parts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text_parts.append(f"[{sheet}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    text_parts.append(' | '.join(cells))
            text_parts.append('')
        wb.close()
        return '\n'.join(text_parts)
    except ImportError:
        raise RuntimeError("请安装 openpyxl: pip install openpyxl")


def _parse_image(path):
    try:
        from ocr_recognizer import get_ocr_instance
        ocr = get_ocr_instance()
        return ocr.recognize(path)
    except Exception as e:
        raise RuntimeError(f"图片 OCR 识别失败: {e}")


def format_file_size(size_bytes):
    if size_bytes < 1024:
        return f"{size_bytes} B"
    elif size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f} KB"
    else:
        return f"{size_bytes / (1024 * 1024):.1f} MB"
